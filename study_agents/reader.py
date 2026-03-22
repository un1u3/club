"""
club/agents/reader.py

Extracts plain text from study materials of various formats.
Acts as the first stage in the CLUB pipeline — raw file to clean text.

Supported formats:
    - .pdf    → PyMuPDF (fitz)
    - .pptx   → python-pptx (all slides)
    - .docx   → python-docx (all paragraphs)
    - .jpg/.jpeg/.png → pytesseract OCR
    - .txt    → direct file read
    - YouTube links   → youtube-transcript-api

Example:
    >>> from study_agents.reader import read_file
    >>> text = read_file("folder/notes/os_chapter5.pdf")
    >>> print(text[:200])

Author: CLUB Project
License: MIT
"""

import os
import re

# ── Constants ─────────────────────────────────────────────

# Supported file extensions mapped to their handler names
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}
TEXT_EXTENSION = ".txt"
PDF_EXTENSION = ".pdf"
PPTX_EXTENSION = ".pptx"
DOCX_EXTENSION = ".docx"

# Path to the file where students paste YouTube links
YOUTUBE_LINKS_FILE = os.path.join(
    "folder", "youtube", "links.txt"
)


# ── Core Function ─────────────────────────────────────────

def read_file(file_path: str) -> str:
    """
    Reads a file and returns its content as clean plain text.

    Detects the file type by extension and delegates to the
    appropriate extraction handler. Strips extra whitespace
    from the result before returning.

    Args:
        file_path: Absolute or relative path to the file.
                   Supports .pdf, .pptx, .docx, .txt, and
                   image files (.jpg, .jpeg, .png).

    Returns:
        Cleaned plain text extracted from the file.
        Returns an empty string if extraction fails.

    Example:
        >>> text = read_file("folder/notes/lecture1.pdf")
        >>> len(text) > 0
        True
    """
    file_extension = os.path.splitext(file_path)[1].lower()

    # Pick the right handler based on file extension
    handler_map = {
        PDF_EXTENSION: _read_pdf,
        PPTX_EXTENSION: _read_pptx,
        DOCX_EXTENSION: _read_docx,
        TEXT_EXTENSION: _read_text,
    }

    if file_extension in handler_map:
        raw_text = handler_map[file_extension](file_path)
    elif file_extension in IMAGE_EXTENSIONS:
        raw_text = _read_image(file_path)
    else:
        print(
            f"CLUB: unsupported file type '{file_extension}' "
            f"for {file_path}"
        )
        return ""

    return _clean_text(raw_text)


# ── Format Handlers ───────────────────────────────────────

def _read_pdf(file_path: str) -> str:
    """
    Extracts text from a PDF using PyMuPDF (fitz).

    Iterates through every page and concatenates extracted
    text with newlines between pages.

    Args:
        file_path: Path to the .pdf file.

    Returns:
        Raw extracted text from all pages, or empty string
        on failure.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print(
            "CLUB: PyMuPDF not installed. "
            "Install with: pip install PyMuPDF"
        )
        return ""

    try:
        document = fitz.open(file_path)
        page_texts = []

        for page in document:
            page_texts.append(page.get_text())

        document.close()
        return "\n".join(page_texts)

    except FileNotFoundError:
        print(f"CLUB: file not found — {file_path}")
        return ""
    except Exception as error:
        print(f"CLUB: failed to read PDF — {error}")
        return ""


def _read_pptx(file_path: str) -> str:
    """
    Extracts text from a PowerPoint file using python-pptx.

    Walks through every slide and every shape, pulling out
    text from each text frame.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        Raw extracted text from all slides, or empty string
        on failure.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "CLUB: python-pptx not installed. "
            "Install with: pip install python-pptx"
        )
        return ""

    try:
        presentation = Presentation(file_path)
        slide_texts = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                # Not all shapes contain text — skip those
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_texts.append(paragraph.text)

        return "\n".join(slide_texts)

    except FileNotFoundError:
        print(f"CLUB: file not found — {file_path}")
        return ""
    except Exception as error:
        print(f"CLUB: failed to read PPTX — {error}")
        return ""


def _read_docx(file_path: str) -> str:
    """
    Extracts text from a Word document using python-docx.

    Reads all paragraphs from the document body and joins
    them with newlines.

    Args:
        file_path: Path to the .docx file.

    Returns:
        Raw extracted text from all paragraphs, or empty
        string on failure.
    """
    try:
        import docx
    except ImportError:
        print(
            "CLUB: python-docx not installed. "
            "Install with: pip install python-docx"
        )
        return ""

    try:
        document = docx.Document(file_path)
        paragraph_texts = [
            paragraph.text for paragraph in document.paragraphs
        ]
        return "\n".join(paragraph_texts)

    except FileNotFoundError:
        print(f"CLUB: file not found — {file_path}")
        return ""
    except Exception as error:
        print(f"CLUB: failed to read DOCX — {error}")
        return ""


def _read_image(file_path: str) -> str:
    """
    Extracts text from an image using pytesseract OCR.

    Opens the image with Pillow and runs Tesseract OCR to
    convert handwritten or printed text into a string.

    Args:
        file_path: Path to the image (.jpg, .jpeg, .png).

    Returns:
        OCR-extracted text, or empty string on failure.

    Raises:
        Prints a warning if Tesseract is not installed on
        the system (separate from the Python package).
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        print(
            "CLUB: pytesseract or Pillow not installed. "
            "Install with: pip install pytesseract pillow"
        )
        return ""

    try:
        image = Image.open(file_path)
        extracted_text = pytesseract.image_to_string(image)
        return extracted_text

    except FileNotFoundError:
        print(f"CLUB: file not found — {file_path}")
        return ""
    except Exception as error:
        # Tesseract binary missing is the most common cause
        print(
            f"CLUB: OCR failed for {file_path} — {error}. "
            "Make sure Tesseract is installed: "
            "sudo apt install tesseract-ocr"
        )
        return ""


def _read_text(file_path: str) -> str:
    """
    Reads a plain text file directly.

    Args:
        file_path: Path to the .txt file.

    Returns:
        File contents as a string, or empty string on
        failure.
    """
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    except FileNotFoundError:
        print(f"CLUB: file not found — {file_path}")
        return ""
    except UnicodeDecodeError:
        print(
            f"CLUB: encoding error reading {file_path} — "
            "file may not be valid UTF-8"
        )
        return ""
    except Exception as error:
        print(f"CLUB: failed to read text file — {error}")
        return ""


def read_youtube_transcripts() -> str:
    """
    Fetches transcripts for all YouTube links in links.txt.

    Reads folder/youtube/links.txt line by line. Each line
    should contain one YouTube URL. Extracts the video ID
    and fetches the transcript via youtube-transcript-api.

    Returns:
        Combined transcript text from all listed videos,
        or empty string if no links or fetch fails.

    Example:
        >>> text = read_youtube_transcripts()
        >>> print(text[:100])
    """
    try:
        from youtube_transcript_api import (
            YouTubeTranscriptApi,
        )
    except ImportError:
        print(
            "CLUB: youtube-transcript-api not installed. "
            "Install with: pip install youtube-transcript-api"
        )
        return ""

    if not os.path.exists(YOUTUBE_LINKS_FILE):
        print(
            f"CLUB: {YOUTUBE_LINKS_FILE} not found — "
            "add YouTube links there, one per line"
        )
        return ""

    with open(YOUTUBE_LINKS_FILE, "r", encoding="utf-8") as file:
        links = file.read().strip().splitlines()

    all_transcripts = []

    for link in links:
        link = link.strip()
        if not link:
            continue

        video_id = _extract_video_id(link)
        if not video_id:
            print(f"CLUB: could not parse video ID from {link}")
            continue

        try:
            transcript_parts = YouTubeTranscriptApi.get_transcript(
                video_id
            )
            # Each part is a dict with 'text', 'start', 'duration'
            transcript_text = " ".join(
                part["text"] for part in transcript_parts
            )
            all_transcripts.append(transcript_text)

        except Exception as error:
            print(
                f"CLUB: failed to fetch transcript for "
                f"{link} — {error}"
            )

    return "\n\n".join(all_transcripts)


# ── Helpers ───────────────────────────────────────────────

def _extract_video_id(url: str) -> str:
    """
    Pulls the YouTube video ID from a URL.

    Handles both formats:
        - https://www.youtube.com/watch?v=VIDEO_ID
        - https://youtu.be/VIDEO_ID

    Args:
        url: A YouTube video URL.

    Returns:
        The 11-character video ID, or empty string if
        the URL format is not recognized.

    Example:
        >>> _extract_video_id("https://youtu.be/dQw4w9WgXcQ")
        'dQw4w9WgXcQ'
    """
    # Pattern covers youtube.com/watch?v= and youtu.be/ links
    pattern = (
        r"(?:youtube\.com/watch\?v=|youtu\.be/)"
        r"([a-zA-Z0-9_-]{11})"
    )
    match = re.search(pattern, url)
    return match.group(1) if match else ""


def _clean_text(raw_text: str) -> str:
    """
    Strips excess whitespace from extracted text.

    Collapses multiple blank lines into one and trims
    leading/trailing whitespace from each line.

    Args:
        raw_text: Unprocessed text from any extractor.

    Returns:
        Cleaned text with normalized whitespace.
    """
    # Collapse runs of 3+ newlines down to 2 (one blank line)
    cleaned = re.sub(r"\n{3,}", "\n\n", raw_text)
    # Strip trailing spaces on each line
    lines = [line.rstrip() for line in cleaned.splitlines()]
    return "\n".join(lines).strip()


# ── Entry Point (Test) ────────────────────────────────────

if __name__ == "__main__":
    import sys

    # Allow passing a file path as a CLI argument for testing
    if len(sys.argv) < 2:
        print(
            "Usage: python -m agents.reader <file_path>\n"
            "Example: python -m agents.reader "
            "folder/notes/lecture1.pdf"
        )
        sys.exit(1)

    test_file_path = sys.argv[1]
    print(f"CLUB Reader: extracting from {test_file_path}\n")

    extracted_text = read_file(test_file_path)

    if extracted_text:
        # Show first 500 chars as a preview
        preview_length = min(500, len(extracted_text))
        print(extracted_text[:preview_length])
        print(f"\n--- ({len(extracted_text)} chars total) ---")
    else:
        print("CLUB: no text extracted.")
